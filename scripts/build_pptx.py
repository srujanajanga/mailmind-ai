#!/usr/bin/env python3
"""Build the MailMind AI presentation deck (docs/MailMind_AI.pptx).

A self-contained python-pptx builder — 14 widescreen slides with a consistent
navy / teal / coral theme, the project's real figures embedded, big-number stat
callouts, and speaker notes. Run:

    PYTHONPATH=src python3 scripts/build_pptx.py
"""
from __future__ import annotations

import sys
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
FIG = ROOT / "docs" / "screenshots"
OUT = ROOT / "docs" / "MailMind_AI.pptx"

# ---- palette -------------------------------------------------------------- #
NAVY = RGBColor(0x1D, 0x35, 0x57)
DARKNAVY = RGBColor(0x12, 0x20, 0x33)
CORAL = RGBColor(0xE6, 0x39, 0x46)
TEAL = RGBColor(0x2A, 0x9D, 0x8F)
SAND = RGBColor(0xE9, 0xC4, 0x6A)
SLATE = RGBColor(0x45, 0x7B, 0x9D)
GRAY = RGBColor(0x6C, 0x75, 0x7D)
INK = RGBColor(0x21, 0x2B, 0x36)
MUTED = RGBColor(0x5B, 0x66, 0x72)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ICE = RGBColor(0xCA, 0xDC, 0xFC)
CARD = RGBColor(0xF4, 0xF6, 0xF9)

TITLE_FONT = "Cambria"
BODY_FONT = "Calibri"

ACCENTS = [CORAL, TEAL, SLATE, SAND, NAVY]

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
W, H = 13.333, 7.5


# --------------------------------------------------------------------------- #
# Low-level helpers
# --------------------------------------------------------------------------- #
def slide(bg=WHITE):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text


def _set_runs(p, text, size, color, font, bold, italic):
    p.text = text
    r = p.runs[0]
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.name = font
    r.font.bold = bold
    r.font.italic = italic


def text(s, txt, x, y, w, h, size=18, color=INK, font=BODY_FONT, bold=False,
         italic=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.04)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    _set_runs(p, txt, size, color, font, bold, italic)
    return tb


def _bullet(p, color):
    """Attach a real round bullet of *color* to paragraph *p*."""
    pPr = p._p.get_or_add_pPr()
    pPr.set("marL", "274320")
    pPr.set("indent", "-274320")
    buClr = pPr.makeelement(qn("a:buClr"), {})
    srgb = pPr.makeelement(qn("a:srgbClr"), {"val": "%02X%02X%02X" % (color[0], color[1], color[2])})
    buClr.append(srgb)
    buFont = pPr.makeelement(qn("a:buFont"), {"typeface": "Arial"})
    buChar = pPr.makeelement(qn("a:buChar"), {"char": "•"})
    for el in (buClr, buFont, buChar):
        pPr.append(el)


def bullets(s, items, x, y, w, h, size=17, color=INK, accent=CORAL, gap=12,
            line=1.12):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.04)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _set_runs(p, item, size, color, BODY_FONT, False, False)
        p.space_after = Pt(gap)
        p.line_spacing = line
        _bullet(p, accent)
    return tb


def rrect(s, x, y, w, h, fill, line_color=None, radius=0.08, shadow=True):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
                             Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(1)
    try:
        shp.adjustments[0] = radius
    except Exception:
        pass
    shp.shadow.inherit = False
    if shadow:
        _soft_shadow(shp)
    return shp


def _soft_shadow(shp):
    spPr = shp._element.spPr
    effLst = spPr.makeelement(qn("a:effectLst"), {})
    outer = spPr.makeelement(qn("a:outerShdw"),
                             {"blurRad": "60000", "dist": "25000", "dir": "5400000",
                              "rotWithShape": "0"})
    clr = spPr.makeelement(qn("a:srgbClr"), {"val": "1D2733"})
    alpha = spPr.makeelement(qn("a:alpha"), {"val": "18000"})
    clr.append(alpha)
    outer.append(clr)
    effLst.append(outer)
    spPr.append(effLst)


def circle(s, x, y, d, fill, transparency=None):
    shp = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    shp.shadow.inherit = False
    if transparency is not None:
        _set_alpha(shp, transparency)
    return shp


def _set_alpha(shp, pct):
    """Set fill transparency (pct 0-100) on a solid-filled shape."""
    srgb = shp.fill.fore_color._xFill.find(qn("a:srgbClr"))
    if srgb is not None:
        a = srgb.makeelement(qn("a:alpha"), {"val": str(int((100 - pct) * 1000))})
        srgb.append(a)


def title_bar(s, kicker, title, accent=CORAL):
    circle(s, 0.7, 0.62, 0.26, accent)
    text(s, kicker.upper(), 1.08, 0.5, 8, 0.35, size=12.5, color=accent,
         font=BODY_FONT, bold=True)
    text(s, title, 1.06, 0.82, 11.5, 0.85, size=30, color=NAVY, font=TITLE_FONT,
         bold=True)


def stat(s, x, y, w, number, label, color, num_size=42, lab_size=13):
    text(s, number, x, y, w, 0.78, size=num_size, color=color, font=TITLE_FONT,
         bold=True, align=PP_ALIGN.LEFT)
    text(s, label, x, y + num_size / 56.0 + 0.18, w, 0.6, size=lab_size,
         color=MUTED, font=BODY_FONT)


def picture_fit(s, path, x, y, max_w, max_h, shadow=True):
    iw, ih = Image.open(path).size
    ar = iw / ih
    w, h = max_w, max_w / ar
    if h > max_h:
        h, w = max_h, max_h * ar
    px = x + (max_w - w) / 2
    py = y + (max_h - h) / 2
    pic = s.shapes.add_picture(str(path), Inches(px), Inches(py), Inches(w), Inches(h))
    if shadow:
        _soft_shadow(pic)
    return pic


def card_block(s, x, y, w, h, heading, body, accent):
    rrect(s, x, y, w, h, CARD)
    circle(s, x + 0.28, y + 0.28, 0.34, accent)
    text(s, heading, x + 0.78, y + 0.22, w - 1.0, 0.5, size=15.5, color=NAVY,
         font=TITLE_FONT, bold=True)
    text(s, body, x + 0.3, y + 0.78, w - 0.6, h - 0.95, size=12.5, color=MUTED,
         font=BODY_FONT)


# =========================================================================== #
# Slide 1 — Title
# =========================================================================== #
s = slide(DARKNAVY)
circle(s, -1.6, -1.6, 4.4, SLATE, transparency=72)
circle(s, W - 2.9, H - 2.9, 4.6, CORAL, transparency=78)
circle(s, W - 1.3, 0.7, 1.3, TEAL, transparency=70)
text(s, "MailMind AI", 1.0, 2.18, 11.3, 1.4, size=64, color=WHITE,
     font=TITLE_FONT, bold=True, align=PP_ALIGN.CENTER)
text(s, "Your Inbox, Intelligently Organized", 1.0, 3.55, 11.3, 0.7, size=27,
     color=SAND, font=TITLE_FONT, italic=True, align=PP_ALIGN.CENTER)
text(s, "NLP  ·  Machine Learning  ·  Agentic AI", 1.0, 4.35, 11.3, 0.5,
     size=16, color=ICE, font=BODY_FONT, align=PP_ALIGN.CENTER)
text(s, "Final-Year Project   |   Team: [Your Names]   |   Guide: [Faculty Guide]",
     1.0, 6.35, 11.3, 0.5, size=13, color=GRAY, font=BODY_FONT,
     align=PP_ALIGN.CENTER)
notes(s, "Open with the one-line promise: an inbox that organises, prioritises "
         "and acts for you. Keep it to ~20 seconds before diving in.")

# =========================================================================== #
# Slide 2 — Problem Statement
# =========================================================================== #
s = slide()
title_bar(s, "The Problem", "Email Overload Is a Productivity Tax", CORAL)
bullets(s, [
    "Knowledge workers face dozens of emails every single day",
    "High cognitive load triaging each message manually",
    "Important and urgent mail gets buried under the noise",
    "Promotions, social and spam crowd out real priorities",
    "Traditional inboxes are flat lists with no sense of importance",
], 0.95, 2.0, 7.4, 4.6, size=18, accent=CORAL, gap=16)
rrect(s, 8.9, 2.1, 3.5, 4.2, CARD)
stat(s, 9.3, 2.55, 2.8, "40+", "emails received per day", CORAL)
stat(s, 9.3, 4.05, 2.8, "~28%", "of the workday spent in the inbox", NAVY)
text(s, "Attention is the scarce resource.", 9.3, 5.45, 2.8, 0.7, size=13,
     color=MUTED, font=BODY_FONT, italic=True)
notes(s, "Frame the pain everyone feels: the inbox is a flat list with no sense "
         "of what matters. Set up the need for intelligence.")

# =========================================================================== #
# Slide 3 — Existing System
# =========================================================================== #
s = slide()
title_bar(s, "Existing System", "Static Rules That Tag, But Don't Understand", SLATE)
cards = [
    ("Manual sorting", "Slow, repetitive and error-prone; the user does all the work.", CORAL),
    ("Keyword / rule filters", "Brittle to new wording; easily fooled and quickly outdated.", SLATE),
    ("No prioritisation", "Cannot rank importance within legitimate mail.", SAND),
    ("No learning", "Ignores how you actually open, reply, ignore or delete.", TEAL),
    ("No assistance", "Produces no summaries and suggests no next action.", NAVY),
]
xs = [0.95, 5.15, 9.35]
for i, (hd, bd, ac) in enumerate(cards[:3]):
    card_block(s, xs[i], 2.1, 3.7, 1.95, hd, bd, ac)
for i, (hd, bd, ac) in enumerate(cards[3:]):
    card_block(s, xs[i], 4.35, 3.7, 1.95, hd, bd, ac)
notes(s, "Stress that traditional filters are static rules — they tag, but they "
         "do not understand, rank or act.")

# =========================================================================== #
# Slide 4 — Proposed System
# =========================================================================== #
s = slide()
title_bar(s, "Proposed System", "MailMind AI — An Agentic Email Assistant", TEAL)
text(s, "Understands, prioritises and acts on your inbox — five pillars in one pipeline.",
     0.95, 1.78, 11.5, 0.5, size=16, color=MUTED, italic=True)
pillars = [
    ("1", "Classification", "6-category ML model", CORAL),
    ("2", "NLP Signals", "keywords · sentiment\nurgency · intent", TEAL),
    ("3", "Context Priority", "sender · recency · VIP", SLATE),
    ("4", "Behaviour", "learns from your actions", SAND),
    ("5", "Agentic Output", "summarise · suggest · flag", NAVY),
]
cw, gap = 2.28, 0.18
x0 = (W - (cw * 5 + gap * 4)) / 2
for i, (n, hd, bd, ac) in enumerate(pillars):
    x = x0 + i * (cw + gap)
    rrect(s, x, 2.55, cw, 3.4, CARD)
    circle(s, x + cw / 2 - 0.38, 2.85, 0.76, ac)
    text(s, n, x + cw / 2 - 0.38, 2.96, 0.76, 0.6, size=26, color=WHITE,
         font=TITLE_FONT, bold=True, align=PP_ALIGN.CENTER)
    text(s, hd, x + 0.1, 3.82, cw - 0.2, 0.6, size=15, color=NAVY,
         font=TITLE_FONT, bold=True, align=PP_ALIGN.CENTER)
    text(s, bd, x + 0.12, 4.42, cw - 0.24, 1.3, size=12, color=MUTED,
         font=BODY_FONT, align=PP_ALIGN.CENTER)
notes(s, "Position MailMind as five pillars working in one pipeline — not just a "
         "smarter filter, but an assistant that takes the next step.")

# =========================================================================== #
# Slide 5 — Objectives
# =========================================================================== #
s = slide()
title_bar(s, "Objectives", "What the System Sets Out to Achieve", SAND)
objs = [
    ("Reduce cognitive load", "Cut the effort of triaging a crowded inbox."),
    ("Prioritise what matters", "Rank mail by true importance and urgency."),
    ("Personalise over time", "Adapt ranking from each user's behaviour."),
    ("Be proactive", "Summaries, urgent flags and suggested actions."),
    ("Improve productivity", "Surface the few emails that need action first."),
]
for i, (hd, bd) in enumerate(objs):
    y = 1.95 + i * 0.96
    circle(s, 0.95, y + 0.06, 0.6, ACCENTS[i])
    text(s, str(i + 1), 0.95, y + 0.12, 0.6, 0.5, size=22, color=WHITE,
         font=TITLE_FONT, bold=True, align=PP_ALIGN.CENTER)
    text(s, hd, 1.8, y, 5.0, 0.55, size=19, color=NAVY, font=TITLE_FONT, bold=True)
    text(s, bd, 1.8, y + 0.45, 10.6, 0.5, size=14, color=MUTED, font=BODY_FONT)
notes(s, "Five objectives map directly to the five pillars; each is measurable "
         "through automation rate and classification accuracy.")

# =========================================================================== #
# Slide 6 — Architecture
# =========================================================================== #
s = slide()
title_bar(s, "Architecture", "A Layered, Modular Python System", SLATE)
picture_fit(s, FIG / "architecture.png", 0.8, 1.75, 11.7, 5.35)
notes(s, "Walk left to right: data in, intelligence in the middle (NLP + ML, then "
         "context + behaviour), agentic insights out to the API and dashboard. A "
         "feedback loop writes user actions back to the store.")

# =========================================================================== #
# Slide 7 — Workflow
# =========================================================================== #
s = slide()
title_bar(s, "Workflow", "How One Email Flows Through MailMind", CORAL)
steps = [
    ("Classify", "ML model assigns one of six categories", CORAL),
    ("Analyse", "NLP extracts keywords, sentiment, urgency, intent", TEAL),
    ("Prioritise", "Context scorer computes a 0–100 priority", SLATE),
    ("Adapt", "Behavioural layer adjusts from past actions", SAND),
    ("Act", "Summarise, suggest actions, flag urgent / VIP / spam", NAVY),
]
for i, (hd, bd, ac) in enumerate(steps):
    y = 1.95 + i * 0.98
    circle(s, 1.0, y, 0.66, ac)
    text(s, str(i + 1), 1.0, y + 0.06, 0.66, 0.55, size=24, color=WHITE,
         font=TITLE_FONT, bold=True, align=PP_ALIGN.CENTER)
    text(s, hd, 1.95, y - 0.04, 3.0, 0.6, size=20, color=NAVY, font=TITLE_FONT, bold=True)
    text(s, bd, 5.0, y + 0.02, 7.6, 0.6, size=15, color=MUTED, font=BODY_FONT)
text(s, "process_inbox() returns insights sorted by priority  ·  record_feedback() closes the loop",
     1.0, 6.95, 11.8, 0.4, size=12.5, color=TEAL, font=BODY_FONT, italic=True)
notes(s, "Emphasise this is a per-email pipeline; record_feedback() closes the "
         "loop so the system adapts over time.")

# =========================================================================== #
# Slide 8 — AI Components
# =========================================================================== #
s = slide()
title_bar(s, "AI Components", "Five Complementary Intelligence Layers", TEAL)
comps = [
    ("Classification", "TF-IDF + engineered features → six categories at 91% accuracy.", CORAL),
    ("NLP Processing", "Keyword extraction, VADER sentiment, urgency and intent detection.", TEAL),
    ("Behavioural Intelligence", "Learns engagement from opens, replies, ignores and deletes.", SLATE),
    ("Contextual Intelligence", "Sender importance, recency and VIP signals → priority.", SAND),
    ("Agentic AI", "Summarises, suggests the next action and adapts over time.", NAVY),
]
for i in range(2):
    for j in range(3):
        k = i * 3 + j
        if k >= len(comps):
            break
        hd, bd, ac = comps[k]
        card_block(s, 0.95 + j * 4.05, 2.05 + i * 2.35, 3.75, 2.05, hd, bd, ac)
notes(s, "One card per pillar. Make clear these are complementary stages, not "
         "competing models.")

# =========================================================================== #
# Slide 9 — Technologies
# =========================================================================== #
s = slide()
title_bar(s, "Technologies", "A Production-Shaped Python Stack", NAVY)
groups = [
    ("Core ML / Data", "scikit-learn · numpy · pandas · scipy · joblib", CORAL),
    ("NLP", "NLTK (VADER, punkt, WordNet) · optional spaCy", TEAL),
    ("Backend / API", "FastAPI · uvicorn · pydantic", SLATE),
    ("Frontend / Viz", "Streamlit · altair · matplotlib · seaborn", SAND),
    ("Storage", "SQLite (standard library)", NAVY),
    ("Testing", "pytest — 44 tests, all passing", TEAL),
]
for i in range(2):
    for j in range(3):
        hd, bd, ac = groups[i * 3 + j]
        x, y = 0.95 + j * 4.05, 2.05 + i * 2.2
        rrect(s, x, y, 3.75, 1.9, CARD)
        circle(s, x + 0.3, y + 0.32, 0.32, ac)
        text(s, hd, x + 0.78, y + 0.26, 2.8, 0.5, size=15.5, color=NAVY,
             font=TITLE_FONT, bold=True)
        text(s, bd, x + 0.32, y + 0.86, 3.1, 0.95, size=13, color=MUTED, font=BODY_FONT)
notes(s, "Group by layer so the audience sees a coherent, production-shaped stack "
         "rather than a random tool list.")

# =========================================================================== #
# Slide 10 — Dataset
# =========================================================================== #
s = slide()
title_bar(s, "Dataset", "A Realistic, Reproducible Synthetic Corpus", SLATE)
bullets(s, [
    "4,200 synthetic emails — 700 per category, perfectly balanced",
    "Six categories with 14 × 14 templates and slot-filling",
    "12% ambiguity borrows confusable neighbour bodies",
    "Category-specific senders, links and attachment rates",
    "Deterministic (seed = 42); stratified 80 / 20 split",
], 0.95, 1.95, 6.4, 4.4, size=16.5, accent=SLATE, gap=14)
text(s, "3,360 train  ·  840 test", 0.95, 6.35, 6.0, 0.5, size=14, color=TEAL,
     font=BODY_FONT, bold=True)
picture_fit(s, FIG / "category_distribution.png", 7.5, 1.95, 5.4, 4.7)
notes(s, "Highlight that the 12% ambiguity makes classes realistically "
         "overlapping — not artificially separable.")

# =========================================================================== #
# Slide 11 — Model Training
# =========================================================================== #
s = slide()
title_bar(s, "Model Training", "Comparing Four Classifiers, Keeping the Best", SAND)
bullets(s, [
    "TF-IDF (1–2)-grams + 8 engineered numeric features",
    "Combined via ColumnTransformer with MaxAbsScaler",
    "Four models compared on the 840-email test set",
    "Best: Logistic Regression — F1 0.9095 / acc 0.9095",
    "Beat Linear SVM, Complement NB and Random Forest",
], 0.95, 1.95, 6.4, 4.4, size=16.5, accent=SAND, gap=14)
picture_fit(s, FIG / "model_comparison.png", 7.4, 1.95, 5.5, 4.7)
notes(s, "Logistic Regression edged out a calibrated Linear SVM and clearly beat "
         "the tree and Naive-Bayes baselines on sparse, high-dimensional text.")

# =========================================================================== #
# Slide 12 — Results & Findings
# =========================================================================== #
s = slide()
title_bar(s, "Results & Findings", "Strong, Honest, Interpretable Performance", CORAL)
rrect(s, 0.95, 1.95, 6.0, 1.55, CARD, shadow=True)
stat(s, 1.25, 2.28, 1.85, "90.9%", "test accuracy", CORAL, num_size=31)
stat(s, 3.35, 2.28, 1.5, "0.91", "macro F1", NAVY, num_size=31)
stat(s, 5.05, 2.28, 1.8, "+18.8", "pts vs rules", TEAL, num_size=31)
bullets(s, [
    "Per-class F1 from 0.898 (Work) to 0.919 (Personal)",
    "Errors cluster on overlapping pairs: Work↔Important, Promotions↔Spam",
    "Rule-based baseline 72.1% → 67.5% relative error reduction",
    "Rules cannot prioritise, learn, summarise or suggest actions",
], 0.95, 3.75, 6.0, 3.2, size=14.5, accent=CORAL, gap=12)
picture_fit(s, FIG / "confusion_matrix.png", 7.2, 1.85, 5.7, 5.2)
notes(s, "The confusion matrix shows errors cluster on genuinely ambiguous pairs, "
         "which is expected and acceptable. The +18.8 points over a keyword "
         "baseline is measured on the same test set.")

# =========================================================================== #
# Slide 13 — Demo / Screenshots
# =========================================================================== #
s = slide()
title_bar(s, "Demo", "The Working Streamlit Dashboard", TEAL)
picture_fit(s, FIG / "dashboard_inbox.png", 0.8, 1.75, 8.7, 5.4)
rrect(s, 9.75, 1.85, 3.05, 5.1, CARD)
text(s, "Live inbox", 10.0, 2.05, 2.6, 0.5, size=16, color=NAVY,
     font=TITLE_FONT, bold=True)
bullets(s, [
    "Sorted by priority score",
    "Critical / High / Medium / Low bands",
    "Category, urgency & intent per email",
    "One-click suggested actions",
    "Outage → Important, High (79.1)",
], 10.0, 2.6, 2.7, 4.0, size=12.5, accent=TEAL, gap=10)
notes(s, "Demo the live inbox — point out the priority ordering and the one-click "
         "suggested actions per email. Analytics and detail views also exist.")

# =========================================================================== #
# Slide 14 — Future Scope & Conclusion
# =========================================================================== #
s = slide(DARKNAVY)
circle(s, W - 3.0, -1.5, 4.2, TEAL, transparency=80)
circle(s, -1.4, H - 2.6, 4.0, CORAL, transparency=82)
text(s, "FUTURE SCOPE & CONCLUSION", 1.0, 0.7, 11, 0.4, size=13, color=SAND,
     font=BODY_FONT, bold=True)
text(s, "From a Flat Inbox to an Intelligent Assistant", 1.0, 1.15, 11.3, 0.9,
     size=30, color=WHITE, font=TITLE_FONT, bold=True)
bullets(s, [
    "Live mailbox integration via IMAP / Gmail & Microsoft Graph",
    "Transformer embeddings (BERT) behind the existing feature interface",
    "Generative AI for context-aware reply drafting",
    "Per-user online learning and multilingual support",
], 1.0, 2.5, 11.4, 2.6, size=18, color=ICE, accent=SAND, gap=16)
text(s, "MailMind AI turns a noisy, flat inbox into a prioritised, self-organising "
        "assistant — classifying, summarising and acting on email with 91% accuracy.",
     1.0, 5.35, 11.4, 1.1, size=16, color=WHITE, font=TITLE_FONT, italic=True)
text(s, "Thank you.", 1.0, 6.55, 11.4, 0.5, size=15, color=SAND, font=BODY_FONT, bold=True)
notes(s, "Be honest — time-saving is an illustrative estimate from the measured "
         "automation rate, not a human user study. Close on the vision.")

# --------------------------------------------------------------------------- #
prs.save(str(OUT))
print(f"Wrote {OUT}  ({len(prs.slides)} slides)")
